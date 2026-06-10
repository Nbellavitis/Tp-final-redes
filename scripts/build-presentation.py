#!/usr/bin/env python3
"""Genera la presentacion final del TPE de Redes (Grupo 12 - Kubernetes)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------- Paleta y constantes ----------
ITBA_BLUE = RGBColor(0x1F, 0x4E, 0x79)        # acento primario
ITBA_BLUE_LIGHT = RGBColor(0x44, 0x72, 0xC4)  # acento secundario
GREY_DARK = RGBColor(0x33, 0x33, 0x33)        # cuerpo
GREY_MID = RGBColor(0x70, 0x70, 0x70)         # subtitulos / footer
GREY_LIGHT = RGBColor(0xE7, 0xEC, 0xF2)       # fondos suaves
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)
ORANGE_WARN = RGBColor(0xC6, 0x6A, 0x00)
RED_BAD = RGBColor(0xB0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

REPO = Path(__file__).resolve().parents[1]
OUT = REPO / "Presentacion-Kubernetes-Grupo12.pptx"
ARCH_IMG = REPO / "docs" / "images" / "architecture.png"


# ---------- Helpers ----------
def add_blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def add_textbox(slide, left, top, width, height, text, *,
                font=FONT, size=18, bold=False, color=GREY_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_round_rect(slide, left, top, width, height, fill, line=None, text=None,
                   text_size=14, text_color=WHITE, bold=True, font=FONT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if text is not None:
        tf = sh.text_frame
        tf.margin_left = Emu(50000)
        tf.margin_right = Emu(50000)
        tf.margin_top = Emu(30000)
        tf.margin_bottom = Emu(30000)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(text_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return sh


def add_line(slide, x1, y1, x2, y2, color=GREY_MID, weight=1.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_title_band(slide, kicker, title, *, accent=ITBA_BLUE):
    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), accent)
    # Kicker
    add_textbox(slide, Inches(0.6), Inches(0.32), Inches(12), Inches(0.35),
                kicker.upper(), size=11, bold=True, color=accent)
    # Title
    add_textbox(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.9),
                title, size=32, bold=True, color=GREY_DARK)
    # Separator
    add_line(slide, Inches(0.6), Inches(1.55), Inches(12.7), Inches(1.55),
             color=GREY_LIGHT, weight=1.5)


def add_footer(slide, page, total):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
                "TPE Redes de Informacion - Grupo 12 - Kubernetes",
                size=9, color=GREY_MID)
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35),
                f"{page} / {total}", size=9, color=GREY_MID, align=PP_ALIGN.RIGHT)


def add_speaker_chip(slide, speaker):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(11.2), Inches(0.3), Inches(2.0), Inches(0.4))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ITBA_BLUE_LIGHT
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Emu(50000); tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = speaker
    run.font.name = FONT; run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = WHITE


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_bullets(slide, left, top, width, height, items, *,
                size=20, color=GREY_DARK, bullet_color=ITBA_BLUE, line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # bullet dot
        dot = p.add_run()
        dot.text = "■  "
        dot.font.name = FONT
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = bullet_color
        # text
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_code_block(slide, left, top, width, height, lines, *, size=12):
    add_rect(slide, left, top, width, height, GREY_LIGHT)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                  width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = GREY_DARK


# ---------- Slides ----------
def slide_title(prs):
    s = add_blank_slide(prs)
    # Left dark band
    add_rect(s, 0, 0, Inches(4.6), SLIDE_H, ITBA_BLUE)
    # Accent stripe
    add_rect(s, Inches(4.6), 0, Inches(0.08), SLIDE_H, ITBA_BLUE_LIGHT)

    add_textbox(s, Inches(0.6), Inches(0.6), Inches(3.8), Inches(0.4),
                "ITBA  -  REDES DE INFORMACION",
                size=11, bold=True, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(6.6), Inches(3.8), Inches(0.4),
                "1er. Cuatrimestre 2026",
                size=11, color=WHITE)

    # Title
    add_textbox(s, Inches(5.1), Inches(2.2), Inches(7.8), Inches(1.2),
                "Despliegue y gestion de un",
                size=36, bold=True, color=GREY_DARK)
    add_textbox(s, Inches(5.1), Inches(2.95), Inches(7.8), Inches(1.2),
                "cluster de Kubernetes",
                size=36, bold=True, color=ITBA_BLUE)

    add_textbox(s, Inches(5.1), Inches(4.1), Inches(7.8), Inches(0.5),
                "POC con Vagrant + Ansible + K3s sobre VMs locales",
                size=18, color=GREY_MID)

    # Separator
    add_line(s, Inches(5.1), Inches(4.85), Inches(8.5), Inches(4.85),
             color=ITBA_BLUE_LIGHT, weight=2.5)

    add_textbox(s, Inches(5.1), Inches(5.0), Inches(7.8), Inches(0.4),
                "Trabajo Practico Especial - Tema 6 - Grupo 12",
                size=14, bold=True, color=GREY_DARK)
    add_textbox(s, Inches(5.1), Inches(5.4), Inches(7.8), Inches(0.4),
                "[Orador 1]  -  [Orador 2]  -  [Orador 3]",
                size=14, color=GREY_MID)
    add_textbox(s, Inches(5.1), Inches(6.6), Inches(7.8), Inches(0.4),
                "Junio 2026",
                size=12, color=GREY_MID)

    add_notes(s,
        "Apertura (Orador 1).\n"
        "Saludar, presentar al grupo y al tema. 30 segundos.\n"
        "Mensaje clave: vamos a mostrar como armar un cluster Kubernetes local "
        "de forma 100% automatica, desplegar una tienda online de microservicios "
        "y demostrar una operacion basica de gestion del cluster.\n"
        "Duracion total objetivo: ~20 min de slides + 8-10 min de demo en vivo.")
    return s


def slide_agenda(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Agenda", "De que vamos a hablar")
    add_speaker_chip(s, "Orador 1")

    items = [
        ("01", "Problematica", "Microservicios, contenedores y por que un cluster"),
        ("02", "Solucion propuesta", "Arquitectura del POC y casos de uso"),
        ("03", "Tecnologias", "Vagrant, Ansible, K3s, Flannel, ingress-nginx"),
        ("04", "Demo en vivo", "Cluster, app, agregar un worker"),
        ("05", "Cierre", "Limites, alternativas y aprendizajes"),
    ]
    top = Inches(2.0)
    for n, title, desc in items:
        add_round_rect(s, Inches(0.6), top, Inches(1.0), Inches(0.8),
                       ITBA_BLUE, text=n, text_size=22)
        add_textbox(s, Inches(1.9), top + Inches(0.05), Inches(10), Inches(0.45),
                    title, size=20, bold=True, color=GREY_DARK)
        add_textbox(s, Inches(1.9), top + Inches(0.45), Inches(10), Inches(0.4),
                    desc, size=14, color=GREY_MID)
        top += Inches(1.0)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Anticipar la estructura: problema, solucion, tecnologias, demo, cierre.\n"
        "Decir explicitamente: 'No hace falta tener conocimiento previo, vamos a "
        "construir los conceptos desde cero'. 30 segundos.")


def slide_aplicacion(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Problematica", "Punto de partida: The Store")
    add_speaker_chip(s, "Orador 1")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(6.5), Inches(0.5),
                "Una tienda online real",
                size=22, bold=True, color=ITBA_BLUE)

    add_bullets(s, Inches(0.6), Inches(2.45), Inches(6.5), Inches(4.0), [
        "5 piezas que cooperan: UI, Catalogo, Carrito, Checkout, Ordenes",
        "Cada pieza es un servicio independiente, escrito en su propio lenguaje",
        "Se comunican por la red, no comparten memoria",
        "Cada una se despliega y escala por separado",
    ], size=16)

    add_textbox(s, Inches(0.6), Inches(5.6), Inches(6.5), Inches(1.2),
                "A esta arquitectura se la llama microservicios.",
                size=14, bold=True, color=GREY_MID)
    add_textbox(s, Inches(0.6), Inches(5.95), Inches(6.5), Inches(1.2),
                "Es flexible y resiliente, pero requiere coordinar muchas piezas.",
                size=14, color=GREY_MID)

    if ARCH_IMG.exists():
        s.shapes.add_picture(str(ARCH_IMG), Inches(7.4), Inches(1.9),
                             width=Inches(5.5))
    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Mostrar el diagrama de la app. Idea: la tienda online no es un solo "
        "programa, son 5 servicios que tienen que hablarse entre si.\n"
        "Analogia util: cada microservicio es como un local distinto del shopping; "
        "el cliente recorre varios locales para terminar la compra.\n"
        "Aclarar: este TP NO es sobre programar la tienda, es sobre como "
        "hacerla correr de forma confiable. ~1 min.")


def slide_concepto_contenedor(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Problematica", "Pieza 1: contenedores")
    add_speaker_chip(s, "Orador 1")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(8), Inches(0.5),
                "Un contenedor empaqueta un servicio con todo lo que necesita",
                size=20, bold=True, color=GREY_DARK)

    # Three boxes
    box_top = Inches(2.7)
    box_h = Inches(2.8)
    box_w = Inches(4.0)
    gap = Inches(0.15)

    # Box 1: codigo
    add_round_rect(s, Inches(0.6), box_top, box_w, Inches(0.5),
                   ITBA_BLUE, text="Codigo", text_size=14)
    add_rect(s, Inches(0.6), box_top + Inches(0.5), box_w, Inches(2.3),
             GREY_LIGHT)
    add_textbox(s, Inches(0.75), box_top + Inches(0.7), box_w - Inches(0.3), Inches(2.0),
                "El programa de la UI, del catalogo, del carrito...",
                size=14, color=GREY_DARK)

    # Box 2: dependencias
    add_round_rect(s, Inches(4.85), box_top, box_w, Inches(0.5),
                   ITBA_BLUE, text="Dependencias", text_size=14)
    add_rect(s, Inches(4.85), box_top + Inches(0.5), box_w, Inches(2.3),
             GREY_LIGHT)
    add_textbox(s, Inches(5.0), box_top + Inches(0.7), box_w - Inches(0.3), Inches(2.0),
                "Librerias, versiones, runtime de Java/Go/Node, sin pisar al host",
                size=14, color=GREY_DARK)

    # Box 3: aislamiento
    add_round_rect(s, Inches(9.1), box_top, box_w, Inches(0.5),
                   ITBA_BLUE, text="Aislamiento", text_size=14)
    add_rect(s, Inches(9.1), box_top + Inches(0.5), box_w, Inches(2.3),
             GREY_LIGHT)
    add_textbox(s, Inches(9.25), box_top + Inches(0.7), box_w - Inches(0.3), Inches(2.0),
                "Procesos, red y archivos separados de otros contenedores",
                size=14, color=GREY_DARK)

    add_textbox(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.5),
                "Idea clave: 'si corre en mi maquina' deja de ser un problema.",
                size=16, bold=True, color=ITBA_BLUE)
    add_textbox(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.5),
                "La imagen del contenedor es la misma en desarrollo, prueba y produccion.",
                size=14, color=GREY_MID)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Definicion simple: un contenedor es una caja con el servicio + todo lo "
        "que necesita para correr, aislado del resto.\n"
        "Analogia: como un envase sellado de comida congelada. Da lo mismo en "
        "que freezer se guarde, sale igual al calentarlo.\n"
        "Importante para el publico no tecnico: NO es una VM, es mucho mas liviano. "
        "Un solo servidor puede correr decenas o cientos.\n"
        "Decir: la herramienta tipica para crear contenedores es Docker. ~1 min.")


def slide_concepto_kubernetes(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Problematica", "Pieza 2: orquestar los contenedores")
    add_speaker_chip(s, "Orador 1")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
                "Tener contenedores no alcanza. Alguien tiene que coordinarlos.",
                size=20, bold=True, color=GREY_DARK)

    # Two columns: "Sin orquestador" vs "Con Kubernetes"
    col_top = Inches(2.7)
    col_h = Inches(3.6)
    col_w = Inches(5.9)

    # Sin orquestador
    add_rect(s, Inches(0.6), col_top, col_w, Inches(0.5), RED_BAD)
    add_textbox(s, Inches(0.6), col_top + Inches(0.07), col_w, Inches(0.4),
                "Sin orquestador", size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.6), col_top + Inches(0.5), col_w, col_h - Inches(0.5),
             GREY_LIGHT)
    add_bullets(s, Inches(0.8), col_top + Inches(0.7), col_w - Inches(0.4),
                col_h - Inches(0.8), [
        "Hay que decirle a cada server que arrancar",
        "Si un contenedor muere, alguien lo levanta a mano",
        "Si hay que escalar, scripts a mano",
        "Las direcciones IP cambian y rompen llamadas",
        "Operacion fragil y no reproducible",
    ], size=13, line_spacing=1.25)

    # Con Kubernetes
    add_rect(s, Inches(6.85), col_top, col_w, Inches(0.5), GREEN_OK)
    add_textbox(s, Inches(6.85), col_top + Inches(0.07), col_w, Inches(0.4),
                "Con Kubernetes", size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.85), col_top + Inches(0.5), col_w, col_h - Inches(0.5),
             GREY_LIGHT)
    add_bullets(s, Inches(7.05), col_top + Inches(0.7), col_w - Inches(0.4),
                col_h - Inches(0.8), [
        "Le decis QUE queres correr (3 replicas de la UI)",
        "El se encarga de DONDE y COMO",
        "Reinicia, reubica y reemplaza solo",
        "Servicios virtuales con nombres estables",
        "Es la base estandar de operacion moderna",
    ], size=13, line_spacing=1.25)

    add_textbox(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
                "Kubernetes = el director de orquesta de los contenedores.",
                size=14, bold=True, color=ITBA_BLUE, align=PP_ALIGN.CENTER)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Idea fuerza: Kubernetes (K8s) es declarativo. Vos describis el estado "
        "deseado y el sistema se ocupa de llegar y mantenerse en ese estado.\n"
        "Analogia: en vez de manejar manualmente cada empleado del shopping, "
        "tenes un gerente que distribuye las tareas y reemplaza al que falto.\n"
        "Es estandar de facto en la industria. ~1 min.")


def slide_problema_real(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Problematica", "El problema real: armar el cluster")
    add_speaker_chip(s, "Orador 1")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
                "Kubernetes resuelve el problema operativo, pero alguien lo tiene que instalar.",
                size=18, color=GREY_DARK)

    # 6 pain points en grilla 3x2
    pains = [
        ("Red mal configurada", "VXLAN por la interfaz incorrecta y los nodos no se ven"),
        ("Puertos cerrados", "El firewall bloquea API server, Flannel o el Ingress"),
        ("Token de union", "Hay que copiarlo a mano de un nodo a otro"),
        ("Versiones distintas", "Cada nodo termina con una version diferente de K8s"),
        ("Sin trazabilidad", "Nadie sabe que cambio en cada nodo ni en que orden"),
        ("Imposible repetir", "Si se rompe, hay que volver a hacer todo a mano"),
    ]
    cols, rows = 3, 2
    cell_w = Inches(4.1)
    cell_h = Inches(1.95)
    start_x = Inches(0.6)
    start_y = Inches(2.6)
    gap = Inches(0.15)
    for idx, (t, d) in enumerate(pains):
        r = idx // cols
        c = idx % cols
        x = start_x + (cell_w + gap) * c
        y = start_y + (cell_h + gap) * r
        add_rect(s, x, y, cell_w, cell_h, GREY_LIGHT)
        add_rect(s, x, y, Inches(0.12), cell_h, RED_BAD)
        add_textbox(s, x + Inches(0.3), y + Inches(0.15), cell_w - Inches(0.5),
                    Inches(0.5), t, size=15, bold=True, color=GREY_DARK)
        add_textbox(s, x + Inches(0.3), y + Inches(0.7), cell_w - Inches(0.5),
                    Inches(1.2), d, size=12, color=GREY_MID)

    add_textbox(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
                "Hacerlo a mano = errores silenciosos, demos que no se repiten y deuda operativa.",
                size=13, bold=True, color=ITBA_BLUE, align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Aterrizar la motivacion del TP: el problema no es 'usar Kubernetes', "
        "es 'armar el cluster donde correra Kubernetes'.\n"
        "Mencionar que este TP, segun la consigna, evalua justamente el "
        "despliegue y la GESTION del cluster, no el desarrollo de la app.\n"
        "Cierre del bloque: pasamos a contar como lo resolvimos. ~1 min.")


# ----- Solucion -----
def slide_propuesta(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Solucion", "Nuestra propuesta")
    add_speaker_chip(s, "Orador 2")

    add_round_rect(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.4),
                   ITBA_BLUE,
                   text="Automatizar la creacion, configuracion y gestion de un\ncluster K3s multi-nodo local, reproducible y auditable.",
                   text_size=22)

    add_textbox(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.5),
                "Tres casos de uso a cubrir",
                size=18, bold=True, color=GREY_DARK)

    use_cases = [
        ("01", "Zero to Kube",
         "Un comando levanta las VMs y deja el cluster Ready"),
        ("02", "Desplegar The Store",
         "La app corre, los servicios se ven y el Ingress responde"),
        ("03", "Agregar un Worker",
         "Sumamos un nodo nuevo sin reinstalar nada"),
    ]
    top = Inches(4.5)
    cell_w = Inches(4.0)
    gap = Inches(0.2)
    for i, (n, t, d) in enumerate(use_cases):
        x = Inches(0.6) + (cell_w + gap) * i
        add_rect(s, x, top, cell_w, Inches(2.0), GREY_LIGHT)
        add_rect(s, x, top, cell_w, Inches(0.08), ITBA_BLUE_LIGHT)
        add_textbox(s, x + Inches(0.25), top + Inches(0.2), cell_w - Inches(0.5),
                    Inches(0.4), n, size=12, bold=True, color=ITBA_BLUE)
        add_textbox(s, x + Inches(0.25), top + Inches(0.55), cell_w - Inches(0.5),
                    Inches(0.5), t, size=16, bold=True, color=GREY_DARK)
        add_textbox(s, x + Inches(0.25), top + Inches(1.05), cell_w - Inches(0.5),
                    Inches(1.0), d, size=12, color=GREY_MID)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Arrancar el bloque de solucion. Mensaje clave: cumplir la pre-entrega "
        "es un contrato. Esto es lo que prometimos.\n"
        "Repasar los 3 casos de uso. La demo los va a tocar a los tres. ~1 min.")


def slide_arquitectura(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Solucion", "Arquitectura del POC")
    add_speaker_chip(s, "Orador 2")

    # Capas: Host -> VMs -> K3s -> The Store
    layer_x = Inches(0.6)
    layer_w = Inches(12.1)
    h = Inches(1.0)
    gap = Inches(0.15)

    # Host
    y0 = Inches(2.0)
    add_round_rect(s, layer_x, y0, layer_w, h, ITBA_BLUE,
                   text="Host (laptop)  -  Vagrant + Ansible + kubectl",
                   text_size=14)

    # VMs
    y1 = y0 + h + gap
    add_round_rect(s, layer_x, y1, layer_w, Inches(0.4), ITBA_BLUE_LIGHT,
                   text="4 VMs Ubuntu Server 24.04 LTS en red privada 192.168.56.0/24",
                   text_size=11)
    # 4 nodos
    node_top = y1 + Inches(0.5)
    node_w = Inches(2.85)
    node_h = Inches(1.5)
    node_gap = Inches(0.1)
    nodes = [
        ("k3s-control", "Control Plane", "192.168.56.10", ITBA_BLUE),
        ("k3s-worker-1", "Worker", "192.168.56.11", GREY_MID),
        ("k3s-worker-2", "Worker", "192.168.56.12", GREY_MID),
        ("k3s-worker-3", "Worker (fase 11)", "192.168.56.13", ORANGE_WARN),
    ]
    for i, (name, role, ip, color) in enumerate(nodes):
        x = layer_x + (node_w + node_gap) * i
        add_rect(s, x, node_top, node_w, node_h, GREY_LIGHT)
        add_rect(s, x, node_top, node_w, Inches(0.08), color)
        add_textbox(s, x + Inches(0.15), node_top + Inches(0.2),
                    node_w - Inches(0.3), Inches(0.4),
                    name, size=13, bold=True, color=GREY_DARK)
        add_textbox(s, x + Inches(0.15), node_top + Inches(0.6),
                    node_w - Inches(0.3), Inches(0.35),
                    role, size=11, color=GREY_MID)
        add_textbox(s, x + Inches(0.15), node_top + Inches(1.0),
                    node_w - Inches(0.3), Inches(0.35),
                    ip, size=11, font=FONT_MONO, color=ITBA_BLUE, bold=True)

    # K3s + ingress
    y2 = node_top + node_h + gap
    add_round_rect(s, layer_x, y2, layer_w, h, ITBA_BLUE,
                   text="K3s v1.35.5  -  Flannel VXLAN  -  ingress-nginx",
                   text_size=14)

    # The Store
    y3 = y2 + h + gap
    add_round_rect(s, layer_x, y3, layer_w, h, GREEN_OK,
                   text="Namespace the-store: UI + Catalog + Carts + Checkout + Orders",
                   text_size=14)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Recorrer la pila de abajo hacia arriba.\n"
        "1) En la laptop tenemos Vagrant, Ansible y kubectl.\n"
        "2) Vagrant crea 4 VMs en una red privada. 3 arrancan, la cuarta queda "
        "apagada por defecto y la levantamos para la demo de gestion.\n"
        "3) Sobre las VMs corre K3s con Flannel para la red entre pods e "
        "ingress-nginx para exponer la UI.\n"
        "4) Arriba de todo, los 5 microservicios de The Store. ~2 min.")


def slide_stack(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Tecnologias", "Stack en 3 herramientas")
    add_speaker_chip(s, "Orador 2")

    tools = [
        ("Vagrant",
         "Define y crea las VMs",
         "Un archivo describe nombre, IP, CPU y RAM de cada nodo. Reproducible y portable.",
         "Infraestructura como codigo"),
        ("Ansible",
         "Configura las VMs",
         "Instala K3s, abre puertos, importa imagenes y aplica los manifiestos. Idempotente.",
         "Automatizacion sin agentes"),
        ("K3s",
         "Es el Kubernetes que corre",
         "Distribucion liviana, certificada y completa. Pensada para edge y para laboratorios.",
         "Orquestacion de contenedores"),
    ]
    top = Inches(2.1)
    cell_w = Inches(4.05)
    cell_h = Inches(4.7)
    gap = Inches(0.15)
    for i, (name, what, how, tag) in enumerate(tools):
        x = Inches(0.6) + (cell_w + gap) * i
        add_rect(s, x, top, cell_w, cell_h, WHITE)
        # header
        add_rect(s, x, top, cell_w, Inches(0.7), ITBA_BLUE)
        add_textbox(s, x, top + Inches(0.15), cell_w, Inches(0.5),
                    name, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # body
        add_rect(s, x, top + Inches(0.7), cell_w, cell_h - Inches(0.7),
                 GREY_LIGHT)
        add_textbox(s, x + Inches(0.25), top + Inches(0.95),
                    cell_w - Inches(0.5), Inches(0.5),
                    what, size=15, bold=True, color=ITBA_BLUE)
        add_textbox(s, x + Inches(0.25), top + Inches(1.55),
                    cell_w - Inches(0.5), Inches(2.5),
                    how, size=13, color=GREY_DARK, line_spacing=1.3)
        add_textbox(s, x + Inches(0.25), top + Inches(4.05),
                    cell_w - Inches(0.5), Inches(0.4),
                    tag.upper(), size=10, bold=True, color=GREY_MID)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Encadenar la metafora: Vagrant pone los ladrillos (las VMs), Ansible "
        "los pinta y los conecta (configuracion), K3s es lo que vive adentro y "
        "hace funcionar la app.\n"
        "Mencionar que las 3 son open source, ampliamente usadas en la industria. ~1 min.")


def slide_vagrant(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Tecnologias", "Vagrant - infraestructura como codigo")
    add_speaker_chip(s, "Orador 2")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(0.5),
                "Como lo usamos", size=20, bold=True, color=ITBA_BLUE)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(4.5), [
        "Box bento/ubuntu-24.04 sobre VirtualBox",
        "Red host-only 192.168.56.0/24 con IPs fijas",
        "3 VMs arrancan automaticamente, la 4ta queda en standby",
        "Definicion declarativa en un solo Ruby file",
        "Levanta todo con: vagrant up",
    ], size=15)

    add_textbox(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.5),
                "Extracto de infra/Vagrantfile",
                size=14, bold=True, color=GREY_MID)
    add_code_block(s, Inches(7.0), Inches(2.4), Inches(5.8), Inches(4.3), [
        "NODES = [",
        "  { name: 'k3s-control',",
        "    ip: '192.168.56.10',",
        "    cpus: 2, memory: 2048 },",
        "  { name: 'k3s-worker-1',",
        "    ip: '192.168.56.11', ... },",
        "  { name: 'k3s-worker-2',",
        "    ip: '192.168.56.12', ... },",
        "  { name: 'k3s-worker-3',",
        "    ip: '192.168.56.13',",
        "    autostart: false }",
        "]",
    ], size=12)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Subrayar: NO clickeamos en VirtualBox. Toda la topologia esta versionada "
        "en Git. Cualquiera del equipo puede recrear el entorno con un comando.\n"
        "La 4ta VM con autostart:false es clave para la demo de fase 11: nos "
        "permite mostrar la topologia inicial 'limpia' y luego sumar el nodo. ~1 min.")


def slide_ansible(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Tecnologias", "Ansible - configuracion automatizada")
    add_speaker_chip(s, "Orador 2")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(12.2), Inches(0.5),
                "Que hace Ansible en este POC, paso por paso",
                size=18, bold=True, color=ITBA_BLUE)

    # Pipeline con 7 etapas
    stages = [
        ("common", "Prepara SO\nkernel, sysctl, firewall"),
        ("k3s_server", "Instala el\nControl Plane"),
        ("k3s_agent", "Une los\nWorkers"),
        ("k3s_kubeconfig", "Exporta el\nkubeconfig"),
        ("ingress", "Despliega\ningress-nginx"),
        ("the_store_images", "Distribuye\nimagenes locales"),
        ("the_store_deploy", "Aplica los\nmanifiestos"),
    ]
    top = Inches(2.7)
    cell_w = Inches(1.7)
    cell_h = Inches(1.7)
    gap = Inches(0.06)
    for i, (name, desc) in enumerate(stages):
        x = Inches(0.6) + (cell_w + gap) * i
        add_rect(s, x, top, cell_w, cell_h, GREY_LIGHT)
        add_rect(s, x, top, cell_w, Inches(0.4), ITBA_BLUE)
        add_textbox(s, x, top + Inches(0.05), cell_w, Inches(0.3),
                    name, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, font=FONT_MONO)
        add_textbox(s, x + Inches(0.1), top + Inches(0.5),
                    cell_w - Inches(0.2), cell_h - Inches(0.5),
                    desc, size=11, color=GREY_DARK, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            arrow_x = x + cell_w + Inches(0.005)
            add_textbox(s, arrow_x - Inches(0.05), top + Inches(0.7),
                        Inches(0.1), Inches(0.4),
                        ">", size=14, bold=True, color=ITBA_BLUE_LIGHT,
                        align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.6), Inches(4.7), Inches(12.2), Inches(0.5),
                "Caracteristicas que nos sirven aca",
                size=16, bold=True, color=ITBA_BLUE)
    add_bullets(s, Inches(0.6), Inches(5.25), Inches(12.2), Inches(1.7), [
        "Idempotencia: re-ejecutar no rompe el entorno; solo cambia lo que difiere",
        "Sin agentes: usa SSH, no hace falta instalar nada en los nodos",
        "Inventario YAML por grupos (control_plane, workers); facil sumar nodos",
        "El token de union K3s lo lee Ansible automaticamente desde el Control Plane",
    ], size=14, line_spacing=1.25)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Recorrer rapido cada rol como un eslabon. Es clave decir que cada rol "
        "es chico y tiene una unica responsabilidad.\n"
        "Idempotencia: si corro el playbook 2 veces seguidas, la segunda no "
        "cambia nada (changed=0). Esto es deseable: nos protege de errores.\n"
        "Token de join: este es el detalle 'pulento' a mostrar. En vez de "
        "copiarlo a mano, Ansible lo lee del Control Plane y lo inyecta en los "
        "workers. Asi automatizamos el caso de uso 3. ~2 min.")


def slide_k3s(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Tecnologias", "K3s - Kubernetes liviano")
    add_speaker_chip(s, "Orador 2")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(0.5),
                "Por que K3s y no Kubernetes 'completo'",
                size=18, bold=True, color=ITBA_BLUE)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(4.5), [
        "Distribucion oficial certificada (CNCF)",
        "Un binario unico < 100 MB",
        "Datastore embebido, no hace falta etcd externo",
        "Pensada para edge, IoT y laboratorios",
        "Las mismas APIs que K8s 'grande': kubectl funciona igual",
    ], size=14)

    add_textbox(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.5),
                "Como esta configurado nuestro K3s",
                size=18, bold=True, color=ITBA_BLUE)
    add_bullets(s, Inches(7.0), Inches(2.4), Inches(5.8), Inches(4.5), [
        "Version pinneada v1.35.5+k3s1",
        "Pod CIDR 10.42.0.0/16  -  Service CIDR 10.43.0.0/16",
        "Flannel VXLAN sobre la interfaz eth1 (red privada)",
        "Traefik desactivado (usamos ingress-nginx)",
        "API Server en 192.168.56.10:6443",
    ], size=14)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Por que pinneamos la version: para que todos los nodos corran exactamente "
        "lo mismo. Evita 'drift'.\n"
        "Por que VXLAN sobre eth1: VirtualBox crea una interfaz NAT (eth0) "
        "duplicada en todas las VMs; si Flannel la elige, la red overlay rompe. "
        "Forzamos eth1 que es la red privada host-only.\n"
        "Por que sacamos Traefik: el manifiesto de The Store referencia "
        "explicitamente la IngressClass 'nginx'. ~1.5 min.")


def slide_red(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Tecnologias", "Red: del host al pod")
    add_speaker_chip(s, "Orador 2")

    # Diagrama: Host -> Ingress (host:80) -> Service -> Pod
    add_textbox(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
                "Como entra una request HTTP a la tienda",
                size=18, bold=True, color=ITBA_BLUE)

    # Boxes con flechas
    elements = [
        ("Host\nhttp://192.168.56.10", "El usuario\n(o curl)", ITBA_BLUE),
        ("ingress-nginx\nen k3s-control", "Recibe en :80\nHost: localhost", ITBA_BLUE_LIGHT),
        ("Service ui\nClusterIP", "Direccion virtual\nestable", ITBA_BLUE_LIGHT),
        ("Pod ui\n10.42.x.x", "El contenedor\nque responde", GREEN_OK),
    ]
    top = Inches(3.0)
    cell_w = Inches(2.85)
    cell_h = Inches(1.6)
    gap = Inches(0.25)
    for i, (head, body, color) in enumerate(elements):
        x = Inches(0.6) + (cell_w + gap) * i
        add_rect(s, x, top, cell_w, cell_h, GREY_LIGHT)
        add_rect(s, x, top, cell_w, Inches(0.55), color)
        add_textbox(s, x, top + Inches(0.05), cell_w, Inches(0.5),
                    head, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), top + Inches(0.7),
                    cell_w - Inches(0.2), cell_h - Inches(0.7),
                    body, size=11, color=GREY_DARK, align=PP_ALIGN.CENTER)
        if i < len(elements) - 1:
            ax = x + cell_w
            add_textbox(s, ax + Inches(0.02), top + Inches(0.6),
                        Inches(0.25), Inches(0.5),
                        ">", size=22, bold=True, color=ITBA_BLUE)

    # Tabla resumen de CIDRs/puertos
    add_textbox(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.4),
                "Rangos y puertos en juego",
                size=15, bold=True, color=ITBA_BLUE)

    headers = ["Capa", "Rango / Puerto", "Quien lo usa"]
    rows = [
        ["Red privada nodos", "192.168.56.0/24",  "SSH, API server, Ingress, VXLAN"],
        ["Pod CIDR",          "10.42.0.0/16",     "IPs internas de cada contenedor"],
        ["Service CIDR",      "10.43.0.0/16",     "Direcciones virtuales de servicios"],
        ["Puertos clave",     "22, 80, 443, 6443, 8472/UDP",
         "SSH, HTTP, HTTPS, API server, Flannel VXLAN"],
    ]
    table_top = Inches(5.55)
    col_widths = [Inches(2.6), Inches(3.4), Inches(6.0)]
    col_x = [Inches(0.6), Inches(0.6) + col_widths[0],
             Inches(0.6) + col_widths[0] + col_widths[1]]
    row_h = Inches(0.35)
    # header
    for cx, cw, txt in zip(col_x, col_widths, headers):
        add_rect(s, cx, table_top, cw, row_h, ITBA_BLUE)
        add_textbox(s, cx + Inches(0.1), table_top + Inches(0.05),
                    cw - Inches(0.2), row_h,
                    txt, size=11, bold=True, color=WHITE)
    for ri, row in enumerate(rows):
        y = table_top + row_h * (ri + 1)
        bg = WHITE if ri % 2 == 0 else GREY_LIGHT
        for cx, cw, txt in zip(col_x, col_widths, row):
            add_rect(s, cx, y, cw, row_h, bg)
            add_textbox(s, cx + Inches(0.1), y + Inches(0.05),
                        cw - Inches(0.2), row_h,
                        txt, size=11, color=GREY_DARK,
                        font=FONT_MONO if "0/" in txt or "/" in txt and txt.count(",") == 0 else FONT)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2.\n"
        "Es el diagrama mas tecnico, pero importante para los docentes.\n"
        "Explicar el flujo: el usuario ataca la IP del Control Plane. "
        "ingress-nginx corre con hostNetwork ahi y publica :80 directo. "
        "Toma la request, mira el header Host y la rutea al Service 'ui'. "
        "El Service es una IP virtual estable del Service CIDR; el balanceo "
        "interno de K3s la traduce a la IP real del pod en el Pod CIDR.\n"
        "Subrayar: ingress-nginx esta SOLO en el Control Plane para tener una "
        "puerta de entrada unica y predecible. ~2 min.")


# ----- Flujo y demo -----
def slide_flujo(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Solucion", "De cero a tienda online en 4 pasos")
    add_speaker_chip(s, "Orador 3")

    steps = [
        ("01", "Build de imagenes",
         "bash scripts/build-images.sh + export-images.sh",
         "Construimos las 5 imagenes Docker y las exportamos a un tar."),
        ("02", "Levantar VMs",
         "cd infra && vagrant up k3s-control k3s-worker-1 k3s-worker-2",
         "Vagrant arranca 3 VMs Ubuntu con IPs fijas."),
        ("03", "Configurar cluster",
         "ansible-playbook playbooks/site.yml",
         "Ansible prepara SO, instala K3s server, une workers, exporta kubeconfig e instala Ingress."),
        ("04", "Desplegar The Store",
         "ansible-playbook playbooks/deploy-store.yml",
         "Importa imagenes en cada nodo y aplica los manifiestos K8s."),
    ]
    top = Inches(2.0)
    row_h = Inches(1.2)
    for i, (n, t, cmd, desc) in enumerate(steps):
        y = top + row_h * i
        add_round_rect(s, Inches(0.6), y + Inches(0.15), Inches(0.9),
                       Inches(0.9), ITBA_BLUE, text=n, text_size=22)
        add_textbox(s, Inches(1.8), y + Inches(0.1), Inches(11), Inches(0.4),
                    t, size=17, bold=True, color=GREY_DARK)
        add_textbox(s, Inches(1.8), y + Inches(0.5), Inches(11), Inches(0.35),
                    cmd, size=11, font=FONT_MONO, color=ITBA_BLUE)
        add_textbox(s, Inches(1.8), y + Inches(0.82), Inches(11), Inches(0.4),
                    desc, size=12, color=GREY_MID)

    add_textbox(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
                "Resultado: cluster Ready + The Store accesible en http://192.168.56.10 (Host: localhost).",
                size=12, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    add_notes(s,
        "Orador 3.\n"
        "Tomar la palabra para mostrar el flujo end-to-end. Mensaje: no son 50 "
        "pasos, son 4 comandos. Cada uno encadena varios roles de Ansible.\n"
        "Que la audiencia se quede con: build -> infra -> config -> deploy. ~1.5 min.")


def slide_validacion(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Solucion", "Como validamos que funciona")
    add_speaker_chip(s, "Orador 3")

    add_textbox(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(0.5),
                "Validacion de infraestructura",
                size=17, bold=True, color=ITBA_BLUE)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(2.0), [
        "Nodos en estado Ready",
        "Pods de The Store en Running",
        "Ingress responde por HTTP",
    ], size=14)

    add_textbox(s, Inches(0.6), Inches(4.4), Inches(6.0), Inches(0.5),
                "Validacion funcional (end-to-end)",
                size=17, bold=True, color=ITBA_BLUE)
    add_bullets(s, Inches(0.6), Inches(4.9), Inches(6.0), Inches(2.4), [
        "scripts/validate-store.sh recorre el flujo de compra real",
        "Home -> catalogo -> producto -> carrito -> checkout -> orden",
        "Comprueba que UI hable con catalog, carts, checkout y orders",
    ], size=14)

    add_textbox(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.5),
                "Comandos de cierre",
                size=17, bold=True, color=ITBA_BLUE)
    add_code_block(s, Inches(7.0), Inches(2.4), Inches(5.8), Inches(4.5), [
        "vagrant status",
        "",
        "kubectl get nodes -o wide",
        "kubectl get pods -n the-store -o wide",
        "kubectl get ingress -n the-store",
        "",
        "curl -H 'Host: localhost' \\",
        "  http://192.168.56.10/",
        "",
        "bash scripts/validate-store.sh",
    ], size=12)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 3.\n"
        "Distincion clave: que un pod este Running no significa que la app sirva. "
        "Por eso ademas tenemos un script de validacion funcional que simula una "
        "compra real recorriendo todos los servicios por el Ingress.\n"
        "Esto cubre el caso de uso 2 de la pre-entrega. ~1 min.")


def slide_gestion(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Solucion", "Gestion basica: agregar un Worker")
    add_speaker_chip(s, "Orador 3")

    add_textbox(s, Inches(0.6), Inches(1.85), Inches(12.2), Inches(0.5),
                "Tercer caso de uso: sumar un nodo sin reinstalar nada",
                size=18, bold=True, color=ITBA_BLUE)

    # Antes / despues
    box_top = Inches(2.7)
    box_w = Inches(5.9)
    box_h = Inches(2.6)
    # Antes
    add_rect(s, Inches(0.6), box_top, box_w, Inches(0.45), GREY_MID)
    add_textbox(s, Inches(0.6), box_top + Inches(0.05), box_w, Inches(0.4),
                "Antes - topologia inicial", size=13, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.6), box_top + Inches(0.45), box_w, box_h - Inches(0.45),
             GREY_LIGHT)
    add_textbox(s, Inches(0.8), box_top + Inches(0.7), box_w - Inches(0.4),
                box_h - Inches(0.7),
                "k3s-control\nk3s-worker-1\nk3s-worker-2",
                size=14, font=FONT_MONO, color=GREY_DARK, line_spacing=1.5)
    # Despues
    add_rect(s, Inches(6.85), box_top, box_w, Inches(0.45), GREEN_OK)
    add_textbox(s, Inches(6.85), box_top + Inches(0.05), box_w, Inches(0.4),
                "Despues - cluster expandido", size=13, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.85), box_top + Inches(0.45), box_w, box_h - Inches(0.45),
             GREY_LIGHT)
    add_textbox(s, Inches(7.05), box_top + Inches(0.7), box_w - Inches(0.4),
                box_h - Inches(0.7),
                "k3s-control\nk3s-worker-1\nk3s-worker-2\nk3s-worker-3  <- nuevo",
                size=14, font=FONT_MONO, color=GREY_DARK, line_spacing=1.5)

    # Pasos
    add_textbox(s, Inches(0.6), Inches(5.55), Inches(12), Inches(0.4),
                "Tres comandos para que el cluster lo integre",
                size=15, bold=True, color=ITBA_BLUE)
    add_code_block(s, Inches(0.6), Inches(6.0), Inches(12.2), Inches(0.9), [
        "vagrant up k3s-worker-3",
        "ansible-playbook playbooks/add-worker.yml -e add_worker_target=k3s-worker-3",
        "ansible-playbook playbooks/add-worker.yml -e import_store_images=true --tags images",
    ], size=11)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 3.\n"
        "Es el momento mas potente del POC: probamos que el cluster se puede "
        "operar en caliente. No reinstalamos nada. Ansible lee el token del "
        "Control Plane, lo inyecta en el nodo nuevo y K3s lo une.\n"
        "Segundo playbook: como las imagenes son locales, hay que importarlas "
        "tambien en el nodo nuevo para evitar ImagePullBackOff cuando el "
        "scheduler ubique pods alli.\n"
        "Resultado: 4 nodos Ready y, tras un rollout, hay pods de carts y orders "
        "corriendo en worker-3. ~1.5 min.")


def slide_demo(prs, page, total):
    s = add_blank_slide(prs)
    # Full color slide
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, ITBA_BLUE)
    add_textbox(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.5),
                "EN VIVO", size=14, bold=True, color=ITBA_BLUE_LIGHT,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(3.0), Inches(12), Inches(1.5),
                "Demo", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_line(s, Inches(5.5), Inches(4.6), Inches(7.8), Inches(4.6),
             color=ITBA_BLUE_LIGHT, weight=2.5)
    add_textbox(s, Inches(0.6), Inches(4.85), Inches(12), Inches(0.5),
                "Cluster operativo  -  The Store funcionando  -  Worker en caliente",
                size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_notes(s,
        "Los tres oradores acompaniando la demo.\n\n"
        "Guion sugerido:\n"
        "1. (cd infra && vagrant status) -> ven las 4 VMs.\n"
        "2. kubectl get nodes -o wide -> ven 3 Ready (sin worker-3 todavia).\n"
        "3. kubectl get pods -n the-store -o wide -> pods Running.\n"
        "4. Abrir http://localhost en navegador con el flujo de compra.\n"
        "5. bash scripts/validate-store.sh -> validacion automatica.\n"
        "6. vagrant up k3s-worker-3 -> levanta la VM extra.\n"
        "7. ansible-playbook add-worker.yml -e add_worker_target=k3s-worker-3.\n"
        "8. kubectl get nodes -o wide -> ahora se ven los 4 Ready.\n"
        "9. kubectl rollout restart deployment -n the-store + get pods -o wide -> "
        "se ven pods schedulados en worker-3.\n\n"
        "Si algo falla en la demo, tener screenshots de respaldo. ~8-10 min.")


# ----- Cierre -----
def slide_alternativas(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Cierre", "Alternativas que consideramos")
    add_speaker_chip(s, "Orador 1")

    alts = [
        ("Minikube",
         "Pensada para 1 o pocos nodos en local.",
         "Oculta la complejidad de bootstrap multi-nodo que queremos mostrar."),
        ("Kind",
         "Kubernetes corriendo dentro de contenedores Docker.",
         "Pierde fidelidad con un escenario real sobre VMs y reduce el rol de Ansible."),
        ("kops o cloud (EKS / GKE)",
         "Productivos, pero requieren cuenta cloud y dependencia de red.",
         "Riesgo de costos y de falla por conectividad durante la exposicion."),
    ]
    top = Inches(2.0)
    cell_h = Inches(1.45)
    gap = Inches(0.15)
    for i, (name, what, why) in enumerate(alts):
        y = top + (cell_h + gap) * i
        add_rect(s, Inches(0.6), y, Inches(12.2), cell_h, GREY_LIGHT)
        add_rect(s, Inches(0.6), y, Inches(0.15), cell_h, ITBA_BLUE)
        add_textbox(s, Inches(0.95), y + Inches(0.15), Inches(11.5), Inches(0.5),
                    name, size=18, bold=True, color=ITBA_BLUE)
        add_textbox(s, Inches(0.95), y + Inches(0.6), Inches(11.5), Inches(0.4),
                    what, size=13, color=GREY_DARK)
        add_textbox(s, Inches(0.95), y + Inches(1.0), Inches(11.5), Inches(0.4),
                    "Por que la descartamos: " + why,
                    size=12, color=GREY_MID)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Mostrar que la decision de K3s + Vagrant + Ansible fue evaluada. "
        "El criterio fue: maximizar el valor pedagogico de mostrar configuracion "
        "de SO, red y union de nodos. Por eso descartamos opciones que lo "
        "abstraen. ~1 min.")


def slide_limites(prs, page, total):
    s = add_blank_slide(prs)
    add_title_band(s, "Cierre", "Limites declarados del POC")
    add_speaker_chip(s, "Orador 1")

    add_textbox(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
                "Por honestidad tecnica: lo que NO esta en alcance",
                size=18, bold=True, color=ITBA_BLUE)

    limits = [
        ("Sin alta disponibilidad del Control Plane",
         "Si cae el nodo control, cae el plano de control. No es el foco de este TP."),
        ("Sin autoscaling",
         "No agregamos nodos automaticamente segun carga; el alta es manual y controlada."),
        ("Sin storage distribuido",
         "Persistencia in-memory. No se evalua perdida de datos ante restart."),
        ("Sin cloud publica",
         "Todo corre local sobre VirtualBox. Sin costo, sin dependencia de red externa."),
    ]
    top = Inches(2.6)
    row_h = Inches(1.0)
    for i, (t, d) in enumerate(limits):
        y = top + row_h * i
        add_round_rect(s, Inches(0.6), y, Inches(0.5), Inches(0.5),
                       ORANGE_WARN, text="!", text_size=18)
        add_textbox(s, Inches(1.3), y, Inches(11.4), Inches(0.45),
                    t, size=16, bold=True, color=GREY_DARK)
        add_textbox(s, Inches(1.3), y + Inches(0.45), Inches(11.4), Inches(0.5),
                    d, size=12, color=GREY_MID)

    add_textbox(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
                "Esto coincide exactamente con lo declarado en la pre-entrega aprobada.",
                size=12, color=GREY_MID, align=PP_ALIGN.CENTER, bold=True)
    add_footer(s, page, total)
    add_notes(s,
        "Orador 1.\n"
        "Decir esto de forma directa. Es buen credito academico mostrar los "
        "limites y que la implementacion respeta el contrato de la pre-entrega.\n"
        "Anticipar posibles preguntas: HA, autoscaling, storage, multi-cloud. "
        "Tener respuestas cortas listas. ~1 min.")


def slide_cierre(prs, page, total):
    s = add_blank_slide(prs)
    # Hero
    add_rect(s, 0, 0, Inches(4.6), SLIDE_H, ITBA_BLUE)
    add_rect(s, Inches(4.6), 0, Inches(0.08), SLIDE_H, ITBA_BLUE_LIGHT)

    add_textbox(s, Inches(0.6), Inches(0.6), Inches(3.8), Inches(0.4),
                "CIERRE", size=11, bold=True, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(2.6), Inches(3.8), Inches(2.5),
                "Lo que nos\nllevamos.",
                size=40, bold=True, color=WHITE, line_spacing=1.0)
    add_textbox(s, Inches(0.6), Inches(6.6), Inches(3.8), Inches(0.4),
                "Grupo 12  -  Tema 6", size=11, color=WHITE)

    add_textbox(s, Inches(5.1), Inches(0.85), Inches(7.8), Inches(0.5),
                "Conclusiones", size=24, bold=True, color=ITBA_BLUE)

    points = [
        ("Infraestructura reproducible",
         "Un Vagrantfile + roles Ansible y cualquiera arma el cluster."),
        ("Automatizacion idempotente",
         "Re-correr el playbook no rompe nada; protege la operacion."),
        ("Cluster Kubernetes real",
         "K3s multi-nodo sobre VMs, no contenedores ni mocks."),
        ("App de microservicios funcional",
         "Validamos el flujo de compra completo, no solo pods Running."),
        ("Gestion en caliente",
         "Sumamos un nodo al cluster sin reinstalar el entorno."),
    ]
    top = Inches(1.5)
    row_h = Inches(1.0)
    for i, (t, d) in enumerate(points):
        y = top + row_h * i
        add_round_rect(s, Inches(5.1), y + Inches(0.05), Inches(0.5),
                       Inches(0.5), ITBA_BLUE_LIGHT,
                       text=str(i + 1), text_size=16)
        add_textbox(s, Inches(5.85), y, Inches(7), Inches(0.45),
                    t, size=16, bold=True, color=GREY_DARK)
        add_textbox(s, Inches(5.85), y + Inches(0.45), Inches(7), Inches(0.5),
                    d, size=12, color=GREY_MID)

    add_footer(s, page, total)
    add_notes(s,
        "Orador 2 + Orador 3 alternados.\n"
        "Cerrar reforzando los 5 puntos. Llamar a preguntas.\n"
        "Tip: dejar abierta la siguiente slide ('Preguntas') al pasarle la "
        "palabra al docente.")


def slide_preguntas(prs, page, total):
    s = add_blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, ITBA_BLUE)
    add_textbox(s, Inches(0.6), Inches(2.8), Inches(12), Inches(1.5),
                "Preguntas",
                size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_line(s, Inches(5.5), Inches(4.4), Inches(7.8), Inches(4.4),
             color=ITBA_BLUE_LIGHT, weight=2.5)
    add_textbox(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
                "Gracias.",
                size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
                "Repo: github / docs/how-to.md, docs/plan-ejecucion-k3s-the-store.md",
                size=11, color=ITBA_BLUE_LIGHT, align=PP_ALIGN.CENTER,
                font=FONT_MONO)
    add_notes(s,
        "Los tres oradores disponibles.\n"
        "Posibles preguntas frecuentes para tener preparadas:\n"
        "- Por que K3s y no K8s 'full'? Por liviano, una distro certificada CNCF.\n"
        "- Por que VXLAN por eth1? Para que Flannel no use la NAT duplicada de VBox.\n"
        "- Como manejarian HA? K3s soporta multi-server con embedded etcd; "
        "queda fuera del POC.\n"
        "- Como pasarian a cloud? Reemplazariamos Vagrant por Terraform y los "
        "playbooks Ansible serian casi los mismos.")


# ---------- Main ----------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,           # 1
        slide_agenda,          # 2
        slide_aplicacion,      # 3
        slide_concepto_contenedor,  # 4
        slide_concepto_kubernetes,  # 5
        slide_problema_real,        # 6
        slide_propuesta,            # 7
        slide_arquitectura,         # 8
        slide_stack,                # 9
        slide_vagrant,              # 10
        slide_ansible,              # 11
        slide_k3s,                  # 12
        slide_red,                  # 13
        slide_flujo,                # 14
        slide_validacion,           # 15
        slide_gestion,              # 16
        slide_demo,                 # 17
        slide_alternativas,         # 18
        slide_limites,              # 19
        slide_cierre,               # 20
        slide_preguntas,            # 21
    ]
    total = len(builders)

    # Slide 1 (titulo) no recibe page/total con footer
    builders[0](prs)
    for i, build in enumerate(builders[1:], start=2):
        build(prs, i, total)

    prs.save(OUT)
    print(f"Generado: {OUT}")


if __name__ == "__main__":
    main()
